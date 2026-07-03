import json
import os
import re
import base64
import tempfile
import urllib.parse
import urllib.request
import ssl
import html
from io import BytesIO
from http.server import BaseHTTPRequestHandler


ALLOWED_EXTENSIONS = {
    '.pdf', '.docx', '.pptx', '.xlsx', '.csv',
    '.jpg', '.jpeg', '.png', '.webp',
}
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.webp'}
MAX_SIZE = 10 * 1024 * 1024


def send_json(start_response, data, status=200):
    status_text = {
        200: 'OK', 400: 'Bad Request', 405: 'Method Not Allowed',
        413: 'Request Entity Too Large', 500: 'Internal Server Error',
    }
    headers = [
        ('Content-Type', 'application/json'),
        ('Access-Control-Allow-Origin', '*'),
    ]
    start_response(f'{status} {status_text.get(status, "Error")}', headers)
    return [json.dumps(data).encode()]


def parse_multipart(environ):
    content_type = environ.get('CONTENT_TYPE', '')
    content_length = int(environ.get('CONTENT_LENGTH', 0) or 0)
    if content_length > MAX_SIZE:
        return None, 'File too large. Maximum size is 10MB.'
    boundary = None
    for part in content_type.split(';'):
        part = part.strip()
        if part.startswith('boundary='):
            boundary = part[9:].strip('"').strip("'")
            break
    if not boundary:
        return None, 'No boundary found in Content-Type'
    body = environ['wsgi.input'].read(content_length)
    boundary_bytes = boundary.encode('utf-8')
    parts = body.split(b'--' + boundary_bytes)
    filename = None
    file_content = None
    for part in parts:
        part = part.strip(b'\r\n').strip(b'\n')
        if part in (b'', b'--', b'\r\n--'):
            continue
        header_end = part.find(b'\r\n\r\n')
        if header_end == -1:
            continue
        headers_raw = part[:header_end].decode('utf-8', errors='replace')
        data = part[header_end:]
        data = re.sub(rb'^[\r\n]+', b'', data)
        data = re.sub(rb'[\r\n]+$', b'', data)
        if 'name="file"' not in headers_raw:
            continue
        fname_match = re.search(r'filename="([^"]*)"', headers_raw)
        if not fname_match:
            continue
        filename = os.path.basename(fname_match.group(1))
        file_content = data
        break
    if not filename or file_content is None:
        return None, 'No file found. Use form field name "file".'
    return (filename, file_content), None


def get_mime_type(ext):
    mime_map = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.csv': 'text/csv',
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.png': 'image/png',
        '.webp': 'image/webp',
    }
    return mime_map.get(ext, 'application/octet-stream')


def describe_image(image_bytes, mime_type):
    api_key = os.environ.get('NVIDIA_API_KEY', '')
    if not api_key:
        return None, 'NVIDIA_API_KEY environment variable not set'
    try:
        b64 = base64.b64encode(image_bytes).decode('utf-8')
        data_uri = f'data:{mime_type};base64,{b64}'
        payload = json.dumps({
            'model': 'meta/llama-3.2-90b-vision-instruct',
            'messages': [{
                'role': 'user',
                'content': [
                    {'type': 'text', 'text': 'Describe this image in detail in Markdown format.'},
                    {'type': 'image_url', 'image_url': {'url': data_uri}},
                ],
            }],
            'max_tokens': 1024,
        }).encode()
        req = urllib.request.Request(
            'https://integrate.api.nvidia.com/v1/chat/completions',
            data=payload,
            headers={'Content-Type': 'application/json', 'Authorization': f'Bearer {api_key}'},
        )
        ctx = ssl.create_default_context()
        with urllib.request.urlopen(req, context=ctx, timeout=30) as resp:
            result = json.loads(resp.read().decode())
            return result['choices'][0]['message']['content'], None
    except Exception as e:
        return None, f'{type(e).__name__}: {e}'


def extract_raw_text(file_path, ext):
    ext = ext.lower()
    try:
        if ext == '.pdf':
            import fitz
            doc = fitz.open(file_path)
            text = '\n'.join(page.get_text() for page in doc)
            doc.close()
            return text
        if ext == '.docx':
            import docx
            doc = docx.Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs)
        if ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text)
            return '\n'.join(texts)
        if ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' '.join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        texts.append(row_text)
            wb.close()
            return '\n'.join(texts)
        if ext == '.csv':
            import csv
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return ' '.join(','.join(row) for row in csv.reader(f))
    except Exception:
        pass
    return None


def strip_markdown(md):
    text = re.sub(r'^#{1,6}\s+', '', md, flags=re.MULTILINE)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    text = re.sub(r'[*_]{1,3}([^*_]+)[*_]{1,3}', r'\1', text)
    text = re.sub(r'`[^`]+`', '', text)
    text = re.sub(r'^\s*[-*+]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*>\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\|', ' ', text)
    text = re.sub(r'^[-*_]{3,}\s*$', '', text, flags=re.MULTILINE)
    text = re.sub(r'^```[\w]*\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'\n{3,}', '\n\n', text).strip()
    return text


def handle_wsgi_request(environ, start_response):
    if environ['REQUEST_METHOD'] != 'POST':
        return send_json(start_response, {'error': 'Method not allowed. Use POST.'}, 405)
    try:
        result, err = parse_multipart(environ)
        if err: return send_json(start_response, {'error': err})
        filename, file_bytes = result
        ext = os.path.splitext(filename)[1].lower()
        if ext not in ALLOWED_EXTENSIONS:
            return send_json(start_response, {'error': f'Unsupported type {ext}'})
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            from markitdown import MarkItDown
            md_engine = MarkItDown()
            result_obj = md_engine.convert(tmp_path)
            markdown_text = result_obj.text_content
            image_desc = None
            image_desc_error = None
            if ext in IMAGE_EXTENSIONS:
                desc, desc_err = describe_image(file_bytes, get_mime_type(ext))
                if desc: markdown_text += f'\n\n## Image Description\n\n{desc}'
                elif desc_err: image_desc_error = desc_err
            import tiktoken
            enc = tiktoken.get_encoding('cl100k_base')
            raw = extract_raw_text(tmp_path, ext) or strip_markdown(markdown_text)
            original_tokens = len(enc.encode(raw)) if raw.strip() else 0
            markdown_tokens = len(enc.encode(markdown_text))
            reduction = round((1 - markdown_tokens / max(original_tokens, 1)) * 100, 1)
            response = {
                'filename': filename,
                'markdown': markdown_text,
                'original_token_count': original_tokens,
                'markdown_token_count': markdown_tokens,
                'reduction_percent': reduction,
            }
            if image_desc_error: response['image_description_error'] = image_desc_error
            return send_json(start_response, response)
        finally:
            if tmp_path and os.path.exists(tmp_path): os.unlink(tmp_path)
    except Exception as e:
        return send_json(start_response, {'error': f'{type(e).__name__}: {e}'})


class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        body = self.rfile.read(int(self.headers.get('content-length', 0)))
        environ = {
            'REQUEST_METHOD': 'POST',
            'CONTENT_TYPE': self.headers.get('content-type', ''),
            'CONTENT_LENGTH': str(len(body)),
            'wsgi.input': BytesIO(body),
        }
        self._run_wsgi(environ)

    def do_GET(self):
        self._send_json({'error': 'Method not allowed. Use POST.'}, 405)

    def _run_wsgi(self, environ):
        status_headers = {}

        def start_response(status, headers):
            status_headers['status'] = status
            status_headers['headers'] = headers

        chunks = handle_wsgi_request(environ, start_response)
        status = int(status_headers.get('status', '200 OK').split()[0])
        headers = dict(status_headers.get('headers', []))
        body = b''.join(chunks)
        self.send_response(status)
        for key, value in headers.items():
            self.send_header(key, value)
        self.end_headers()
        self.wfile.write(body)

    def _send_json(self, data, status=200):
        body = json.dumps(data).encode()
        self.send_response(status)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(body)
